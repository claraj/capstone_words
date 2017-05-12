import html2text
from os import listdir, path
import re
import html2text
from pptx import Presentation


def get_all_words():
    word_source_dir = 'words_from_capstone'
    all_text_file = 'allwords.txt'

    all_files = listdir(word_source_dir)

    out_lines = []

    for f in all_files:


        filepath = path.join(word_source_dir, f)
        print(filepath)

        if f.endswith('ppt') or f.endswith('pptx'):
            text = read_ppt(filepath)
            out_lines.append(text)

        elif f.endswith('html'):
            text = read_html(filepath)
            out_lines.append(text)

        elif f.endswith('xml'):
            text = read_xml(filepath)
            out_lines.append(text)


    with open(all_text_file, 'w') as f:
        for line in out_lines:
            f.write(line)
        f.close()





def read_xml(filename):

    with open(filename) as xml_file:
        text = xml_file.read()

        # remove xml tags e.g. <content>
        # remove encoded HTML tags e.g <P> encoded as &lt;p&gt;

        tag_re ='<.+?>';
        encoded_tag_re = '&lt;.+?&gt;'  # ? for ungreedy

        text = re.sub(tag_re, '', text)
        text = re.sub(encoded_tag_re, '', text)

        text = text.replace('nbsp;', '')   # Ugh
        text = text.replace('amp;', '')    # There's probably a better way to read this XML

        return text


def read_ppt(filename):

    text_runs = []
    pres = Presentation(filename)

    for slide in pres.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text_runs.append(run.text)

    return '\n'.join(text_runs)


def read_html(filename):

    with open(filename) as html_file:
        html = html_file.read()
        h = html2text.HTML2Text()
        h.ignore_links = True
        text = h.handle(html)


        return text


get_all_words()
